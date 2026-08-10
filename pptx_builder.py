"""
Builds a 'Player Presentation' scouting-report deck from analyzed play data,
matching the look of the DefensiveIQ sample deck (now rebranded ScoutEdge
Defense): dark title/red-zone slides,
white content slides with metric cards and colored data tables.
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from analyzer import ZONE_ORDER, ZONE_LABEL, MIN_SAMPLE_FORMATION, top_n_counts, run_pass_split

# ------------------------------------------------------------------
# Palette / type — pulled directly from the sample DefensiveIQ deck
# ------------------------------------------------------------------
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_FILL = RGBColor(0xF2, 0xF2, 0xF2)
LABEL_GRAY = RGBColor(0x55, 0x55, 0x55)
BODY_TEXT = RGBColor(0x11, 0x11, 0x11)
RUN_RED = RGBColor(0x8B, 0x00, 0x00)
PASS_BLUE = RGBColor(0x00, 0x00, 0x8B)
EXPL_RUN = RGBColor(0xC0, 0x39, 0x2B)
EXPL_PASS = RGBColor(0x1A, 0x52, 0x76)
TEAL = RGBColor(0x0E, 0x70, 0x60)
LOGO_BORDER = RGBColor(0x3A, 0x4A, 0x6A)
FOOTER_GRAY = RGBColor(0x6A, 0x7A, 0x9A)
SUBTITLE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
RZ_LABEL_FILL = RGBColor(0xC0, 0x39, 0x2B)
RZ_VALUE_FILL = RGBColor(0x22, 0x30, 0x50)

TITLE_FONT = "Cambria"
BODY_FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _bg(slide, color):
    rect = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)  # MSO_SHAPE.RECTANGLE == 1
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    rect.shadow.inherit = False
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)
    return rect


def _emu(v):
    """Round any EMU-ish value (which may have become a plain float after
    arithmetic like width/2) down to a clean int, since python-pptx requires
    integer EMUs and silently writes invalid XML otherwise."""
    return int(round(v))


def _textbox(slide, left, top, width, height, text, size, color, bold=False, italic=False,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(_emu(left), _emu(top), _emu(width), _emu(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    return box


def _title_bar(slide, text, color=BLACK):
    _textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), text, 28, color, bold=True, font=TITLE_FONT)


def build_title_slide(prs, opponent, week, game_date, plays_analyzed):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, BLACK)

    for x in (Inches(1.0), Inches(9.9)):
        box = slide.shapes.add_shape(1, x, Inches(0.6), Inches(2.4), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = BLACK
        box.line.color.rgb = LOGO_BORDER
        box.line.width = Pt(0.75)
        box.shadow.inherit = False
    _textbox(slide, Inches(1.0), Inches(0.6), Inches(2.4), Inches(1.3), "[ YOUR LOGO ]", 11,
             FOOTER_GRAY, font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _textbox(slide, Inches(9.9), Inches(0.6), Inches(2.4), Inches(1.3), "[ OPP LOGO ]", 11,
             FOOTER_GRAY, font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _textbox(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2), "OPPONENT SCOUTING REPORT", 44,
             WHITE, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.9), opponent.upper(), 32,
             WHITE, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
    date_str = game_date.strftime("%B %d, %Y") if hasattr(game_date, "strftime") else str(game_date)
    _textbox(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.5),
             f"Week {week}  ·  {date_str}  ·  {plays_analyzed} plays analyzed", 16,
             SUBTITLE_BLUE, font=BODY_FONT, align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4),
             "ScoutEdge Defense  ·  Auto-generated from film — replace logos and decorate freely",
             10, FOOTER_GRAY, italic=True, font=BODY_FONT, align=PP_ALIGN.CENTER)
    return slide


def _metric_card(slide, left, top, width, height, number, label, number_color):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_FILL
    card.line.fill.background()
    card.shadow.inherit = False
    _textbox(slide, left, top + Inches(0.25), width, Inches(0.9), number, 48, number_color,
             bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _textbox(slide, left, top + height - Inches(0.4), width, Inches(0.4), label, 12, LABEL_GRAY,
             bold=True, font=BODY_FONT, align=PP_ALIGN.CENTER)


def build_overview_slide(prs, df, tendencies):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _title_bar(slide, "OFFENSIVE OVERVIEW")

    runs = (df["PLAY TYPE"] == "Run").sum()
    passes = (df["PLAY TYPE"] == "Pass").sum()
    total = runs + passes
    expl_run = df[df["PLAY TYPE"] == "Run"]["EXPLOSIVE"].mean() if runs else 0
    expl_pass = df[df["PLAY TYPE"] == "Pass"]["EXPLOSIVE"].mean() if passes else 0

    metrics = [
        (str(total), "TOTAL PLAYS", BLACK),
        (f"{round(runs/total*100) if total else 0}%", "RUN %", RUN_RED),
        (f"{round(passes/total*100) if total else 0}%", "PASS %", PASS_BLUE),
        (f"{round(expl_run*100)}%", "EXPL RUN %", EXPL_RUN),
        (f"{round(expl_pass*100)}%", "EXPL PASS %", EXPL_PASS),
    ]
    card_w, gap, left0, top = Inches(2.3), Inches(0.15), Inches(0.6), Inches(1.5)
    for i, (num, label, color) in enumerate(metrics):
        left = left0 + i * (card_w + gap)
        _metric_card(slide, left, top, card_w, Inches(1.8), num, label, color)

    _textbox(slide, Inches(0.6), Inches(3.7), Inches(11), Inches(0.5), "BIGGEST TENDENCIES", 20,
             EXPL_RUN, bold=True, font=BODY_FONT)
    y = Inches(4.4)
    for t in tendencies:
        sq = slide.shapes.add_shape(1, Inches(0.7), y + Inches(0.05), Inches(0.22), Inches(0.22))
        sq.fill.solid()
        sq.fill.fore_color.rgb = EXPL_RUN
        sq.line.fill.background()
        sq.shadow.inherit = False
        _textbox(slide, Inches(1.1), y - Inches(0.05), Inches(10.5), Inches(0.4), t, 15,
                 BODY_TEXT, font=BODY_FONT)
        y += Inches(0.55)
    return slide


def _add_table(slide, left, top, width, height, header, rows, col_widths,
                header_fill, colored_cols=None, label_col_fill=None):
    """colored_cols: dict {col_index: RGBColor} to color that column's text bold in body rows."""
    n_rows, n_cols = len(rows) + 1, len(header)
    gframe = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gframe.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for c, htext in enumerate(header):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        run = p.add_run()
        run.text = htext
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.name = BODY_FONT
        run.font.color.rgb = WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for r, rowvals in enumerate(rows, start=1):
        for c, val in enumerate(rowvals):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = label_col_fill if (c == 0 and label_col_fill) else CARD_FILL
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "" if val is None else str(val)
            run.font.size = Pt(11 if c == 0 else 10.5)
            font_color = WHITE if (c == 0 and label_col_fill) else BODY_TEXT
            if colored_cols and c in colored_cols:
                font_color = colored_cols[c]
                run.font.bold = True
            elif c == 0:
                run.font.bold = True
            run.font.name = BODY_FONT
            run.font.color.rgb = font_color
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return gframe


DD_SLIDE_ROWS = [
    ("P & 10", "P & 10"),
    ("1st & 10", "1ST & 10"), ("2nd & Short", "2ND & SHORT"), ("2nd & Medium", "2ND & MEDIUM"),
    ("2nd & Long", "2ND & LONG"), ("3rd & Short", "3RD & SHORT"), ("3rd & Medium", "3RD & MEDIUM"),
    ("3rd & Long", "3RD & LONG"), ("4th Down", "4TH DOWN"),
]


def build_down_distance_slide(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _title_bar(slide, "DOWN & DISTANCE")

    header = ["SITUATION", "RUN%", "PASS%", "TOP RUN", "TOP PASS", "TOP FORM"]
    rows = []
    for label, bucket in DD_SLIDE_ROWS:
        sub = df[df["DD_BUCKET"] == bucket]
        runs, passes, total, rp, pp = run_pass_split(sub)
        top_run = top_n_counts(sub[sub["PLAY TYPE"] == "Run"]["CONCEPT"], 1)[0].split(" (")[0]
        top_pass = top_n_counts(sub[sub["PLAY TYPE"] == "Pass"]["CONCEPT"], 1)[0].split(" (")[0]
        top_form = top_n_counts(sub["FORMATION"], 1)[0].split(" (")[0]
        rows.append([label, f"{round(rp*100) if rp is not None else 0}%",
                     f"{round(pp*100) if pp is not None else 0}%", top_run, top_pass, top_form])

    col_widths = [Inches(2.0), Inches(1.2), Inches(1.2), Inches(2.7), Inches(2.7), Inches(2.3)]
    tbl = _add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.8), header, rows, col_widths,
                      header_fill=BLACK, colored_cols={1: RUN_RED, 2: PASS_BLUE}, label_col_fill=TEAL)
    return slide


def build_concepts_slide(prs, df, play_type, title, expl_note, header_fill, last_col):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _title_bar(slide, title)
    _textbox(slide, Inches(0.6), Inches(1.05), Inches(11), Inches(0.35), expl_note, 12, LABEL_GRAY, italic=True)

    sub = df[df["PLAY TYPE"] == play_type]
    counts = sub["CONCEPT"].value_counts()
    header = [f"{play_type.upper()} CONCEPT", "CALLED", "AVG YD", "EXPL%", "SUCC%", last_col]
    rows = []
    for concept, called in counts.items():
        if called < 3 or pd.isna(concept):
            continue
        cdf = sub[sub["CONCEPT"] == concept]
        avg_yd = cdf["GAIN"].mean()
        expl_pct = cdf["EXPLOSIVE"].mean()
        succ = cdf["SUCCESS"].dropna()
        succ_pct = succ.mean() if len(succ) else None
        extra = top_n_counts(cdf["FORMATION"], 1)[0].split(" (")[0] if last_col == "TOP FORM" else "—"
        rows.append([
            concept, int(called), round(avg_yd, 1) if pd.notna(avg_yd) else "—",
            f"{round(expl_pct*100)}%", f"{round(succ_pct*100)}%" if succ_pct is not None else "—", extra,
        ])
    col_widths = [Inches(3.0), Inches(1.5), Inches(1.5), Inches(1.5), Inches(1.5), Inches(3.1)]
    _add_table(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(4.2), header, rows, col_widths,
               header_fill=header_fill)
    return slide


def build_formation_slide(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, WHITE)
    _title_bar(slide, "FORMATION TENDENCIES")

    header = ["FORMATION", "SNAPS", "RUN%", "PASS%", "TOP RUN", "TOP PASS"]
    rows = []
    for form, snaps in df["FORMATION"].value_counts().items():
        if snaps < MIN_SAMPLE_FORMATION or pd.isna(form):
            continue
        fdf = df[df["FORMATION"] == form]
        runs, passes, total, rp, pp = run_pass_split(fdf)
        top_run = top_n_counts(fdf[fdf["PLAY TYPE"] == "Run"]["CONCEPT"], 1)[0].split(" (")[0]
        top_pass = top_n_counts(fdf[fdf["PLAY TYPE"] == "Pass"]["CONCEPT"], 1)[0].split(" (")[0]
        rows.append([form, int(snaps), f"{round(rp*100)}%", f"{round(pp*100)}%", top_run, top_pass])
    col_widths = [Inches(2.4), Inches(1.4), Inches(1.4), Inches(1.4), Inches(2.7), Inches(2.8)]
    _add_table(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(4.2), header, rows, col_widths,
               header_fill=EXPL_PASS)
    return slide


def _kv_row(slide, y, label, value):
    lbl = slide.shapes.add_shape(1, Inches(0.8), y, Inches(3.4), Inches(0.6))
    lbl.fill.solid(); lbl.fill.fore_color.rgb = RZ_LABEL_FILL; lbl.line.fill.background(); lbl.shadow.inherit = False
    _textbox(slide, Inches(0.95), y, Inches(3.1), Inches(0.6), label, 15, WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE)
    val = slide.shapes.add_shape(1, Inches(4.3), y, Inches(8.2), Inches(0.6))
    val.fill.solid(); val.fill.fore_color.rgb = RZ_VALUE_FILL; val.line.fill.background(); val.shadow.inherit = False
    _textbox(slide, Inches(4.45), y, Inches(7.9), Inches(0.6), value, 15, WHITE, anchor=MSO_ANCHOR.MIDDLE)


def build_red_zone_slide(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, BLACK)
    _textbox(slide, Inches(0.6), Inches(0.4), Inches(11), Inches(0.8), "RED ZONE KEYS", 28, WHITE,
             bold=True, font=TITLE_FONT)

    rzdf = df[df["ZONE"] == "RZ"]
    gldf = df[df["ZONE"] == "GL"]
    runs, passes, total, rp, pp = run_pass_split(rzdf)
    top_run = top_n_counts(rzdf[rzdf["PLAY TYPE"] == "Run"]["CONCEPT"], 1)[0].split(" (")[0]
    top_pass = top_n_counts(rzdf[rzdf["PLAY TYPE"] == "Pass"]["CONCEPT"], 1)[0].split(" (")[0]
    top_form = top_n_counts(rzdf["FORMATION"], 1)[0].split(" (")[0]
    gl_top = top_n_counts(gldf["CONCEPT"], 1)[0].split(" (")[0] if len(gldf) else "—"

    rows = [
        ("Inside 20", f"Run {round((rp or 0)*100)}%  /  Pass {round((pp or 0)*100)}%   ({total} plays)"),
        ("Top Run", top_run),
        ("Top Pass", top_pass),
        ("Top Formation", top_form),
        ("Top Personnel", "Not tagged"),
        ("Goal Line Favorite", gl_top),
    ]
    y = Inches(1.7)
    for label, value in rows:
        _kv_row(slide, y, label, value)
        y += Inches(0.72)

    _textbox(slide, Inches(0.6), Inches(6.9), Inches(11), Inches(0.4),
              "Small samples possible in RZ/GL — verify against film.", 11, FOOTER_GRAY, italic=True)
    return slide


def biggest_tendencies(df, n=5):
    candidates = []
    for form, cnt in df["FORMATION"].value_counts().items():
        if cnt < MIN_SAMPLE_FORMATION or pd.isna(form):
            continue
        fdf = df[df["FORMATION"] == form]
        runs, passes, total, rp, pp = run_pass_split(fdf)
        skew = max(rp, pp)
        tag = "Run" if rp >= pp else "Pass"
        flag = " *" if cnt < 5 else ""
        candidates.append((skew, f"{form} = {round(skew*100)}% {tag}{flag}", cnt))
    for label, bucket in DD_SLIDE_ROWS:
        sub = df[df["DD_BUCKET"] == bucket]
        if len(sub) < 3:
            continue
        runs, passes, total, rp, pp = run_pass_split(sub)
        skew = max(rp, pp)
        tag = "Run" if rp >= pp else "Pass"
        candidates.append((skew, f"{round(skew*100)}% {tag} on {label}", total))
    candidates.sort(key=lambda x: -x[0])
    out = []
    seen = set()
    for skew, text, total in candidates:
        if text in seen:
            continue
        seen.add(text)
        out.append(f"{text}   ({total} plays)")
        if len(out) >= n:
            break
    return out


def build_presentation(df, opponent="Opponent", week="1", game_date=None):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total_plays = len(df)
    tendencies = biggest_tendencies(df)

    build_title_slide(prs, opponent, week, game_date or "", total_plays)
    build_overview_slide(prs, df, tendencies)
    build_down_distance_slide(prs, df)
    build_concepts_slide(prs, df, "Run", "FAVORITE RUNS",
                         "Explosive = 10+ yards.  Success by down-adjusted yardage.", RUN_RED, "TOP FORM")
    build_concepts_slide(prs, df, "Pass", "FAVORITE PASSES",
                         "Explosive = 15+ yards.", PASS_BLUE, "TOP FORM")
    build_formation_slide(prs, df)
    build_red_zone_slide(prs, df)

    # Lazy import avoids a circular import (hit_chart_builder imports shared
    # colors/helpers from this module).
    from hit_chart_builder import build_hit_chart_slides
    build_hit_chart_slides(prs, df, opponent=opponent)

    from situation_hit_chart_builder import build_situation_hit_chart_slides
    build_situation_hit_chart_slides(prs, df, opponent=opponent)

    return prs
