"""
Down & Distance Hit Charts — same 4-panel-per-slide format as the Formation
Hit Charts, but keyed by situation (1st & 10, 2nd & Short, 2nd & Medium,
2nd & Long, 3rd & Short, 3rd & Medium, 3rd & Long, 4th Down) instead of by
formation. Since a down/distance situation doesn't have one alignment to
draw, the alignment diagram is replaced with a run/pass proportion bar and
a "top formations" line — everything else (top run/pass concepts box) is
shared with the formation hit charts.
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from analyzer import top_n_counts, run_pass_split
from pptx_builder import (
    BLACK, WHITE, BODY_TEXT, RUN_RED, PASS_BLUE, LABEL_GRAY,
    TITLE_FONT, BODY_FONT, SLIDE_W, SLIDE_H, _bg, _textbox, _emu,
)
from hit_chart_builder import draw_run_pass_split_box

TEAL = (0x0E, 0x70, 0x60)
PURPLE = (0x4A, 0x23, 0x5A)
MAROON = (0x7B, 0x24, 0x1C)

DD_SITUATIONS_ORDERED = [
    ("1st & 10", "1ST & 10"),
    ("1st & Short", "1ST & SHORT"),
    ("2nd & Short", "2ND & SHORT"),
    ("2nd & Medium", "2ND & MEDIUM"),
    ("2nd & Long", "2ND & LONG"),
    ("3rd & Short", "3RD & SHORT"),
    ("3rd & Medium", "3RD & MEDIUM"),
    ("3rd & Long", "3RD & LONG"),
    ("4th Down", "4TH DOWN"),
]


def situation_color(bucket):
    """Down-progression urgency color, matching the same scheme used in the
    Excel workbook: teal (1st, safe) -> blue (2nd) -> red (3rd) -> maroon
    (4th, most critical)."""
    s = bucket.upper()
    if "4TH" in s:
        return MAROON
    if "3RD" in s:
        return RUN_RED
    if "2ND" in s:
        return PASS_BLUE
    return TEAL


def _proportion_bar(slide, left, top, width, height, run_pct, pass_pct):
    run_pct = run_pct or 0
    pass_pct = pass_pct or 0
    total_pct = run_pct + pass_pct
    if total_pct <= 0:
        run_pct, pass_pct = 0.5, 0.5
    run_w = width * run_pct
    pass_w = width - run_w

    if run_w > Inches(0.05):
        run_seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(left), _emu(top), _emu(run_w), _emu(height))
        run_seg.fill.solid(); run_seg.fill.fore_color.rgb = RGBColor(0x8B, 0x00, 0x00)
        run_seg.line.fill.background(); run_seg.shadow.inherit = False
    if pass_w > Inches(0.05):
        pass_seg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(left + run_w), _emu(top), _emu(pass_w), _emu(height))
        pass_seg.fill.solid(); pass_seg.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x8B)
        pass_seg.line.fill.background(); pass_seg.shadow.inherit = False

    if run_w > Inches(0.5):
        _textbox(slide, left, top, run_w, height, f"{round(run_pct*100)}% R", 11, WHITE,
                  bold=True, font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if pass_w > Inches(0.5):
        _textbox(slide, left + run_w, top, pass_w, height, f"{round(pass_pct*100)}% P", 11, WHITE,
                  bold=True, font=BODY_FONT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _situation_panel(slide, left, top, width, height, label, bucket, sub):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(left), _emu(top), _emu(width), _emu(height))
    box.fill.background()
    box.line.color.rgb = BLACK
    box.line.width = Pt(1.25)
    box.shadow.inherit = False

    count = len(sub)
    runs, passes, total, rp, pp = run_pass_split(sub)
    color = situation_color(bucket)

    _textbox(slide, left + Inches(0.12), top + Inches(0.06), width - Inches(0.24), Inches(0.3),
              f"{label}  ({count}x)", 14, RGBColor(*color), bold=True, font=BODY_FONT)

    bar_top = top + Inches(0.42)
    _proportion_bar(slide, left + Inches(0.12), bar_top, width - Inches(0.24), Inches(0.3), rp, pp)

    forms = top_n_counts(sub["FORMATION"], 3)
    forms_text = "Top Formations: " + (", ".join(f for f in forms if f != "—") or "—")
    _textbox(slide, left + Inches(0.12), bar_top + Inches(0.38), width - Inches(0.24), Inches(0.26),
              forms_text, 10.5, BLACK, font=BODY_FONT)

    split_top = bar_top + Inches(0.70)
    split_h = top + height - split_top - Inches(0.08)
    run_concepts = top_n_counts(sub[sub["PLAY TYPE"] == "Run"]["CONCEPT"], 3)
    pass_concepts = top_n_counts(sub[sub["PLAY TYPE"] == "Pass"]["CONCEPT"], 3)
    draw_run_pass_split_box(slide, left + Inches(0.1), split_top, width - Inches(0.2), split_h,
                             run_concepts, pass_concepts)


def build_situation_hit_chart_slides(prs, df, opponent="Opponent", min_snaps=1):
    situations = [(label, bucket) for label, bucket in DD_SITUATIONS_ORDERED
                  if len(df[df["DD_BUCKET"] == bucket]) >= min_snaps]
    if not situations:
        return prs

    chunks = [situations[i:i + 4] for i in range(0, len(situations), 4)]
    positions = [
        (Inches(0.4), Inches(1.2), Inches(6.1), Inches(2.9)),
        (Inches(6.75), Inches(1.2), Inches(6.1), Inches(2.9)),
        (Inches(0.4), Inches(4.25), Inches(6.1), Inches(2.9)),
        (Inches(6.75), Inches(4.25), Inches(6.1), Inches(2.9)),
    ]
    for chunk in chunks:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, WHITE)
        _textbox(slide, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.5),
                  "DOWN & DISTANCE HIT CHART", 24, BLACK, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
        _textbox(slide, Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.35),
                  f"DATA FROM: {opponent.upper()}", 13, LABEL_GRAY, italic=True, font=BODY_FONT, align=PP_ALIGN.CENTER)
        for (left, top, w, h), (label, bucket) in zip(positions, chunk):
            sub = df[df["DD_BUCKET"] == bucket]
            _situation_panel(slide, left, top, w, h, label, bucket, sub)
    return prs
