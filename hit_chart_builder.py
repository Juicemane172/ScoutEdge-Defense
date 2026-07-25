"""
Formation Hit Charts — the classic coaching-staff format: one panel per
formation showing snap count, run/pass split, backfield tendency, a quick
alignment diagram, and the top run/pass concepts out of it. Four panels
per slide, matching the layout coaches actually use on the sideline.

LABEL CONVENTION (per DeAirus's own terminology, not the raw Hudl export):
  T = running back, always.
  H = tight end (11 personnel) OR the slot/hybrid receiver in 10 personnel
      (1 RB, no TE) — same alignment spot, different player type.
  F = second running back (20 personnel only).
  X = the boundary / short-side receiver — always drawn on the left.
  Z = the field / wide-side receiver — always drawn on the right.
  A = an extra detached slot receiver on the field/strong side, when the
      formation has one (e.g. Trips).
  Q = quarterback.

HONEST NOTE: which side is actually "strong" varies play to play with the
hash and call; there's no single fixed left/right in real film. This
diagram draws the strong/receiver-heavy side on the right (the Z side) as
a representative, standard convention — like the "Trips Right" example —
not a claim about which literal side this opponent lines up on most often.
Personnel (11/10/20) is read straight from your PERSONNEL column, so the
TE-vs-no-TE and one-back-vs-two-back calls are real; the exact letter
placement of individual slot receivers is the generic part.
"""

import pandas as pd
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from analyzer import top_n_counts, run_pass_split
from pptx_builder import (
    BLACK, WHITE, CARD_FILL, BODY_TEXT, RUN_RED, PASS_BLUE, LABEL_GRAY,
    TITLE_FONT, BODY_FONT, SLIDE_W, SLIDE_H, _bg, _textbox, _emu,
)

HIT_CHART_MIN_SNAPS = 7
HIT_CHART_TOP_N = 6

# Base formation names that are inherently spread/open sets (per DeAirus's own
# naming convention) — in these, even an 11-personnel "TE" is flexed out at
# receiver depth/spacing, not attached to the line. Anything else (Pro, Trey,
# Doubles, etc.) is a TE-based formation where the TE is genuinely attached —
# either on the line, or off the ball if the name carries a "YO" tag.
OPEN_BASE_NAMES = {"DEUCE", "TRIPS", "SLOT", "BUNCH", "CLUSTER"}

# (weak_side_total, strong_side_total) detached-receiver counts per side,
# including any attached H/TE that flexes out as part of that side's count.
FAMILY_SPLIT = {"1X1": (1, 1), "2X1": (1, 2), "2X2": (2, 2), "3X1": (1, 3), "3X2": (2, 3)}

TE_COLOR = RUN_RED
RB_COLOR = (0xE8, 0xC4, 0x2A)  # yellow, matches the sample hit chart
WR_COLOR = PASS_BLUE
QB_COLOR = (0x11, 0x11, 0x11)

# (te_count, rb_count, wr_count) by PERSONNEL grouping
PERSONNEL_MAP = {11.0: (1, 1, 3), 20.0: (0, 2, 3), 10.0: (0, 1, 4)}

DOT_R = Inches(0.095)


def select_hit_chart_formations(df, min_snaps=HIT_CHART_MIN_SNAPS, top_n=HIT_CHART_TOP_N):
    out = []
    for form, cnt in df["FORMATION"].value_counts().items():
        if pd.isna(form) or cnt < min_snaps:
            continue
        out.append(form)
        if len(out) >= top_n:
            break
    return out


def classify_formation_tags(formation_name):
    """Reads DeAirus's own formation-naming tags straight off the FORMATION
    column: a 'YO' tag means the TE is off the ball; a 'FIB' tag (Formation
    In Boundary) means this formation's numbers strength is set toward the
    boundary instead of the field. Base name (first token) decides whether
    the formation is an inherently open/spread set at all."""
    name = (formation_name or "").upper().replace("-", " ")
    tokens = name.split()
    base = tokens[0] if tokens else ""
    is_open_base = base in OPEN_BASE_NAMES
    has_yo = "YO" in tokens
    is_fib = "FIB" in tokens

    if is_open_base:
        te_style = "open"
    elif has_yo:
        te_style = "off_ball"
    else:
        te_style = "inline"
    return te_style, is_fib


def _resolve_personnel(personnel):
    """Returns (te, rb, wr_detached, draw_h, h_is_te) per the T/H convention above."""
    te, rb, wr = PERSONNEL_MAP.get(personnel, (1, 1, 3))
    if te >= 1:
        return te, rb, wr, True, True        # 11 personnel: H = TE
    if rb == 1:
        return te, rb, max(wr - 1, 2), True, False   # 10 personnel: H = slot receiver
    return te, rb, wr, False, False          # 20 personnel (or other, no TE): no H drawn


def _split_extras(family, wr_detached, draw_h):
    """How many 'extra' detached receivers (beyond the fixed X/Z outside pair)
    go on the weak side vs the strong side, based on the formation's actual
    FORMATION FAMILY (2x2, 3x1, etc.) — this is what makes Deuce actually
    look like a 2x2 and Trips actually look like a 3x1. (Which physical side
    -- left/boundary or right/field -- "strong" maps to is decided
    separately, from the FIB tag.)"""
    weak_total, strong_total = FAMILY_SPLIT.get(family, (2, 2))
    weak_share = max(0, weak_total - 1)               # minus the X spot
    strong_share = max(0, strong_total - 1 - (1 if draw_h else 0))  # minus Z, minus H
    total_extra = max(0, wr_detached - 2)              # minus X and Z
    total_share = weak_share + strong_share
    if total_share == 0:
        # Family split doesn't call for any inside receivers (e.g. 1x1), but
        # personnel data has one anyway — default it to the strong side.
        return 0, total_extra
    weak_extra = round(total_extra * weak_share / total_share)
    strong_extra = total_extra - weak_extra
    return weak_extra, strong_extra


def _dot(slide, cx, cy, r, color, label, label_color=WHITE):
    from pptx.dml.color import RGBColor
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, _emu(cx - r), _emu(cy - r), _emu(r * 2), _emu(r * 2))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*color) if isinstance(color, tuple) else color
    shp.line.color.rgb = RGBColor(0x33, 0x33, 0x33)
    shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(8)
    run.font.bold = True
    run.font.name = BODY_FONT
    run.font.color.rgb = RGBColor(*label_color) if isinstance(label_color, tuple) else label_color


def draw_alignment_diagram(slide, left, top, width, height, family, personnel, te_style="inline", is_fib=False):
    """Draws a compact, non-overlapping formation alignment diagram inside the given box.

    te_style: 'open' (TE flexed out at WR depth/spacing), 'inline' (TE tight
      to the tackle, on the line), or 'off_ball' (TE tight to the tackle
      horizontally, stepped back off the line — a "Y-off" look).
    is_fib: if True, the formation's numbers-strength is set toward the
      boundary (left/X side) instead of the default field (right/Z) side.
    """
    te, rb, wr_detached, draw_h, h_is_te = _resolve_personnel(personnel)
    r = DOT_R

    cx = left + width / 2
    los_y = top + height * 0.34

    # sign = +1 means the strong side is on the right (field/Z, the default);
    # sign = -1 means FIB has flipped it to the left (boundary/X).
    sign = -1 if is_fib else 1

    # --- Offensive line: 2 dots, black center square, 2 dots ---
    ol_gap = Inches(0.30)
    for dx in (-2, -1, 1, 2):
        _dot(slide, cx + dx * ol_gap, los_y, r, (0xFF, 0xFF, 0xFF), "", label_color=(0, 0, 0))
    from pptx.dml.color import RGBColor
    center_sq = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(cx - r), _emu(los_y - r), _emu(r * 2), _emu(r * 2))
    center_sq.fill.solid(); center_sq.fill.fore_color.rgb = RGBColor(0, 0, 0)
    center_sq.line.fill.background(); center_sq.shadow.inherit = False

    # --- QB and backfield, staggered depth so nothing stacks ---
    qb_y = los_y + Inches(0.26)
    _dot(slide, cx, qb_y, r, QB_COLOR, "Q")
    rb_y = qb_y + Inches(0.26)
    if rb >= 1:
        _dot(slide, cx + Inches(0.16), rb_y, r, RB_COLOR, "T", label_color=(0, 0, 0))
    if rb >= 2:
        _dot(slide, cx - Inches(0.16), rb_y, r, RB_COLOR, "F", label_color=(0, 0, 0))

    # --- Fixed edges: X is always boundary (left), Z is always field (right) ---
    left_edge = left + Inches(0.18)
    right_edge = left + width - Inches(0.18)
    margin = Inches(0.30)
    _dot(slide, left_edge, los_y, r, WR_COLOR, "X")
    _dot(slide, right_edge, los_y, r, WR_COLOR, "Z")

    # Strong/weak side physical positions, mirrored by `sign` for FIB.
    strong_edge = right_edge if sign == 1 else left_edge
    weak_edge = left_edge if sign == 1 else right_edge
    strong_tackle_x = cx + sign * 2 * ol_gap
    weak_tackle_x = cx - sign * 2 * ol_gap

    weak_extra, strong_extra = _split_extras(family, wr_detached, draw_h)
    inline_gap = Inches(0.20)     # H sitting tight against the tackle
    off_ball_depth = Inches(0.16)  # Y-off: same horizontal spot, stepped back
    slot_step = Inches(0.30)

    def interp(start_x, end_x, frac):
        return start_x + (end_x - start_x) * frac

    if draw_h and te_style in ("inline", "off_ball"):
        # H is attached tight to the strong-side tackle; strong extras (A, Y)
        # start spacing out beyond H.
        h_x = strong_tackle_x + sign * inline_gap
        h_y = los_y if te_style == "inline" else los_y + off_ball_depth
        h_color = TE_COLOR if h_is_te else WR_COLOR
        _dot(slide, h_x, h_y, r, h_color, "H")
        strong_start_x = h_x + sign * slot_step
        strong_end_x = strong_edge - sign * margin
        strong_labels = ["A", "Y"]
        for i in range(min(strong_extra, 2)):
            frac = (i + 1) / (min(strong_extra, 2) + 1)
            _dot(slide, interp(strong_start_x, strong_end_x, frac), los_y, r, WR_COLOR, strong_labels[i])
    else:
        # 'open' TE (or no TE at all): H, if present, joins the WR spacing
        # progression as the innermost strong-side receiver instead of
        # sitting tight to the line.
        strong_items = []
        if draw_h and te_style == "open":
            strong_items.append(("H", TE_COLOR if h_is_te else WR_COLOR))
        strong_items += [(lbl, WR_COLOR) for lbl in ["A", "Y"][:strong_extra]]
        strong_start_x = strong_tackle_x + sign * slot_step
        strong_end_x = strong_edge - sign * margin
        n = len(strong_items)
        for i, (label, color) in enumerate(strong_items):
            frac = (i + 1) / (n + 1)
            _dot(slide, interp(strong_start_x, strong_end_x, frac), los_y, r, color, label)

    # Weak side: mirror the same near-to-far spacing (no H ever routes here).
    weak_start_x = weak_tackle_x - sign * slot_step
    weak_end_x = weak_edge + sign * margin
    weak_labels = ["B", "C"]
    for i in range(min(weak_extra, 2)):
        frac = (i + 1) / (min(weak_extra, 2) + 1)
        _dot(slide, interp(weak_start_x, weak_end_x, frac), los_y, r, WR_COLOR, weak_labels[i])


def _panel(slide, left, top, width, height, form, fdf):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(left), _emu(top), _emu(width), _emu(height))
    box.fill.background()
    box.line.color.rgb = BLACK
    box.line.width = Pt(1.25)
    box.shadow.inherit = False

    count = len(fdf)
    runs, passes, total, rp, pp = run_pass_split(fdf)
    _textbox(slide, left + Inches(0.12), top + Inches(0.06), width * 0.65, Inches(0.3),
              f"Formation: {form} ({count}x)", 13, BLACK, bold=True, font=BODY_FONT)
    rp_text = f"R/P: {round((rp or 0)*100)}%R/{round((pp or 0)*100)}%P"
    _textbox(slide, left + width * 0.62, top + Inches(0.06), width * 0.36, Inches(0.3),
              rp_text, 12, BLACK, bold=True, font=BODY_FONT, align=PP_ALIGN.RIGHT)

    backfields = top_n_counts(fdf["BACKFIELD"], 2)
    bf_parts = []
    for entry in backfields:
        if entry == "—":
            continue
        name, cnt_str = entry.rsplit(" (", 1)
        cnt = int(cnt_str.rstrip(")"))
        pct = round(cnt / count * 100) if count else 0
        bf_parts.append(f"{name}({pct}%)")
    _textbox(slide, left + Inches(0.12), top + Inches(0.36), width - Inches(0.24), Inches(0.28),
              "Backfield: " + "  ".join(bf_parts) if bf_parts else "Backfield: —", 11.5, BLACK,
              bold=True, font=BODY_FONT)

    personnel = fdf["PERSONNEL"].mode()
    personnel = personnel.iloc[0] if not personnel.empty else 11.0
    family = fdf["FORMATION FAMILY"].mode() if "FORMATION FAMILY" in fdf.columns else pd.Series([])
    family = family.iloc[0] if not family.empty else "2X2"
    te_style, is_fib = classify_formation_tags(form)
    diagram_top = top + Inches(0.66)
    diagram_h = height * 0.42
    draw_alignment_diagram(slide, left, diagram_top, width, diagram_h, family, personnel,
                            te_style=te_style, is_fib=is_fib)

    # Top runs / top passes split box
    split_top = diagram_top + diagram_h + Inches(0.05)
    split_h = top + height - split_top - Inches(0.08)
    split_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _emu(left + Inches(0.1)), _emu(split_top),
                                        _emu(width - Inches(0.2)), _emu(split_h))
    split_box.fill.background()
    split_box.line.color.rgb = BLACK
    split_box.line.width = Pt(1)
    split_box.shadow.inherit = False
    mid_x = left + width / 2
    divider = slide.shapes.add_connector(1, _emu(mid_x), _emu(split_top), _emu(mid_x), _emu(split_top + split_h))
    divider.line.color.rgb = BLACK
    divider.line.width = Pt(1)

    run_concepts = top_n_counts(fdf[fdf["PLAY TYPE"] == "Run"]["CONCEPT"], 3)
    pass_concepts = top_n_counts(fdf[fdf["PLAY TYPE"] == "Pass"]["CONCEPT"], 3)

    _textbox(slide, left + Inches(0.16), split_top + Inches(0.04), width / 2 - Inches(0.24), Inches(0.2),
              "TOP RUNS:", 10.5, BLACK, bold=True, font=BODY_FONT)
    _textbox(slide, mid_x + Inches(0.1), split_top + Inches(0.04), width / 2 - Inches(0.24), Inches(0.2),
              "TOP PASSES:", 10.5, BLACK, bold=True, font=BODY_FONT)

    def bullets(items):
        lines = []
        for entry in items:
            if entry == "—":
                continue
            name, cnt_str = entry.rsplit(" (", 1)
            cnt = cnt_str.rstrip(")")
            lines.append(f"\u2022  {name} - {cnt}X")
        return "\n".join(lines) if lines else "\u2022  —"

    run_box = slide.shapes.add_textbox(_emu(left + Inches(0.16)), _emu(split_top + Inches(0.26)),
                                        _emu(width / 2 - Inches(0.3)), _emu(split_h - Inches(0.3)))
    tf = run_box.text_frame; tf.word_wrap = True
    for i, line in enumerate(bullets(run_concepts).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line; r.font.size = Pt(10); r.font.name = BODY_FONT; r.font.color.rgb = BODY_TEXT

    pass_box = slide.shapes.add_textbox(_emu(mid_x + Inches(0.1)), _emu(split_top + Inches(0.26)),
                                        _emu(width / 2 - Inches(0.3)), _emu(split_h - Inches(0.3)))
    tf = pass_box.text_frame; tf.word_wrap = True
    for i, line in enumerate(bullets(pass_concepts).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = line; r.font.size = Pt(10); r.font.name = BODY_FONT; r.font.color.rgb = BODY_TEXT


def build_hit_chart_slides(prs, df, opponent="Opponent"):
    formations = select_hit_chart_formations(df)
    if not formations:
        return prs

    chunks = [formations[i:i + 4] for i in range(0, len(formations), 4)]
    positions = [
        (Inches(0.4), Inches(1.2), Inches(6.1), Inches(2.9)),
        (Inches(6.75), Inches(1.2), Inches(6.1), Inches(2.9)),
        (Inches(0.4), Inches(4.25), Inches(6.1), Inches(2.9)),
        (Inches(6.75), Inches(4.25), Inches(6.1), Inches(2.9)),
    ]
    for chunk in chunks:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, WHITE)
        _textbox(slide, Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.5),
                  "TOP FORMATIONS HIT CHART", 24, BLACK, bold=True, font=TITLE_FONT, align=PP_ALIGN.CENTER)
        _textbox(slide, Inches(0.4), Inches(0.72), Inches(12.5), Inches(0.35),
                  f"DATA FROM: {opponent.upper()}", 13, LABEL_GRAY, italic=True, font=BODY_FONT, align=PP_ALIGN.CENTER)
        for (left, top, w, h), form in zip(positions, chunk):
            fdf = df[df["FORMATION"] == form]
            _panel(slide, left, top, w, h, form, fdf)
    return prs
